from typing import List, Dict, Any, Optional
import fitz
import docx
import re
import io
from pptx import Presentation
import os
from langchain.text_splitter import NLTKTextSplitter, RecursiveCharacterTextSplitter
import nltk
from nltk.tokenize import word_tokenize
from rank_bm25 import BM25Okapi
import requests
import faiss
import numpy as np
from app.core.config import config
from app.db.crud import add_vectors, search_vector
from openai import OpenAI

from app.db.database import FaissSingleton
faiss_client = FaissSingleton()
text_list = []

nltk.download('punkt')
nltk.download('wordnet')
nltk.download('words')
nltk.download('punkt_tab')
os.environ['TESSDATA_PREFIX'] = '/usr/share/tesseract-ocr/4.00/tessdata/'
os.environ['OPENAI_API_KEY'] = config.OPENAI_API_KEY

client = OpenAI(api_key=config.OPENAI_API_KEY)
text_splitter_nltk = NLTKTextSplitter(chunk_size=256, chunk_overlap=50)
text_splitter_recursive = RecursiveCharacterTextSplitter(chunk_size=256*4, chunk_overlap=50)

system_prompt = '''
You are a Question Answering and Text Highlighting expert specialized in extracting and highlighting relevant quotes from texts to answer specific questions\
in a single sentence. If the text is included in an answer and is quoted from the provided text then it must be enclosed between *** without any change in words. 

Following are examples of questions and the text that would contain the answer to the question:
Example 1: 
Question: How can social media impact love relationships negatively?
	Text: Commenting and liking other people's posts and direct messages
	Good output: Social media can impact love relationships negatively when people are ***commenting and liking other people's posts and direct messages***.
	Bad output: Social media can impact love relationships negatively when people comment and like other people's posts and direct messages.
    Bad output: Social media can impact love relationships negatively when people have ***commented and liked other people's posts and direct messages*** (the word forms are changed).
    Bad output: Social media can impact love relationships negatively when ***people are commenting and liking other people's posts and direct messages***. (extra words are included in the text between ***)

Example 2:
Question: What impact could rent caps have on the housing market?
	Text: They argued that a cap on rents would lead landlords to sell their rental properties to owner occupants
	Good output: Rent caps could impact the housing market by leading ***landlords to sell their rental properties to owner occupants***.
	Bad output: Rent caps could impact the housing market by leading landlords to sell their rental properties to owner occupants.
    Bad output: Rent caps could impact the housing market by ***leading landlords to sell their rental properties to owner occupants*** (the word forms are changed).
    
IMPORTANT INSTRUCTIONS:
1. Your answer should be a single sentence and should be relevant to the question and should not contain any carrage returns or new line characters.
2. Your answer should include a direct quote from the text without any change in words or word forms which is enclosed between *** and not double quotation marks.
3. The answer should not contain any new lines and must be a single sentence.
4. The quote between *** should be an exact excerpt from the original text without any change in the words, punctuation or capitalization as the text excerpt needs to be searchable with a Ctrl + F command.
5. Do not in any case change the words of the provided text rather craft your answer in a way that it includes an original excerpt from the text.
6. The answer should not contain anything else other than the answer with the quote from the text.
7. The answer should not contain any additional information or context.

Only the text that is in the original text must be included in the answer 
For Example:
Text: All vaccines teach the immune system to create antibodies to help it fight off a particular pathogen.
Bad Answer: Vaccines work in the body by ***teaching the immune system to create antibodies to help it fight off a particular pathogen.***
Good Answer: Vaccines work in the body by teaching ***the immune system to create antibodies to help it fight off a particular pathogen.***
As illustrated in the example above, the word "teaching" is not included in the highlighted answer as it is not in the original text.

The users question and text will be in the following format:
QUESTION: the question
TEXT: the text
'''

def clean_text(text):
    # Remove consecutive newlines
    cleaned_text = re.sub(r'\n+', ' ', text)
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text)
    cleaned_text = re.sub(r'\xa0', ' ', cleaned_text)
    cleaned_text = re.sub(r'[\u2022\u2023\u25E6\u2043\u2219\u25AA]', '', cleaned_text)
    return cleaned_text.strip()

def word_count(text):
    words = text.split()
    return len(words)

def page_has_image(page):
    for block in page.get_text("blocks"):
        if block[4] == 4:  # 4 means this block is an image
            return True
    return False

def extract_text_from_txt(file_content: bytes) -> str:
    text = file_content.decode('latin-1')  # Decode bytes to string using latin-1 encoding
    cleaned_text = clean_text(text)  # Clean the whole text
    return cleaned_text

def extract_text_from_pdf(file_content: bytes) -> list:
    with fitz.open(stream=file_content, filetype="pdf") as pdf_document:
        text = []
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            blocks = page.get_text("blocks")
            if len(blocks) == 1 and blocks[0][4] == 4:  # check if the page contains only a single image block
                print(f"Page {page_num} is an image. Performing OCR...")
                pix = blocks[0][-1]  # get Pixmap object from image block
                # OCR the image, make a 1-page PDF from it
                pdfdata = pix.pdfocr_tobytes()  # 1-page PDF in memory
                ocrpdf = fitz.open("pdf", pdfdata)  # open as PDF document
                ocrtext = ocrpdf[0].get_text()
                text.append({"page":page_num, "text":ocrtext})
            else:
                text.append({"page":page_num, "text":page.get_text()})
        return text
    
def extract_whole_text_from_pdf(file_content: bytes) -> str:
    with fitz.open(stream=file_content, filetype="pdf") as pdf_document:
        full_text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            blocks = page.get_text("blocks")
            if len(blocks) == 1 and blocks[0][4] == 4:  # Check if the page contains only a single image block
                print(f"Page {page_num} is an image. Performing OCR...")
                pix = blocks[0][-1]  # Get Pixmap object from image block
                # OCR the image, make a 1-page PDF from it
                pdfdata = pix.pdfocr_tobytes()  # 1-page PDF in memory
                ocrpdf = fitz.open("pdf", pdfdata)  # Open as PDF document
                ocrtext = ocrpdf[0].get_text()
                full_text += ocrtext
            else:
                full_text += page.get_text()
        return (clean_text(full_text))

def extract_text_from_docx(file_content: bytes) -> str:
    document = docx.Document(io.BytesIO(file_content))
    full_text = []
    
    for para in document.paragraphs:
        if para.text.strip():  # Only include non-empty paragraphs
            full_text.append(clean_text(para.text))
    
    return "\n".join(full_text)

def extract_text_from_pptx(file_content: bytes) -> str:
    prs = Presentation(io.BytesIO(file_content))
    slide_text = []
    
    for slide in prs.slides:
        slide_text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text_content.append(shape.text)
        if slide_text_content:  # Only include non-empty slides
            slide_text.append(clean_text("\n".join(slide_text_content)))
    
    return "\n".join(slide_text) 


def get_chunks(text: str, chunk_size: int = 128) -> List[str]:
    try:
        final_chunks = []
        chunks = text_splitter_nltk.split_text(text)
        for chunk in chunks:
            if len(word_tokenize(chunk)) > chunk_size:
                final_chunks.extend(text_splitter_recursive.split_text(chunk))
            else:
                final_chunks.append(chunk)
        global text_list
        text_list.extend(final_chunks)
        print (text_list)
        return final_chunks
    except Exception as e:
        print ("*****************ALARM at get_chunks*****************")
        return str(e)

async def process_documents_async(text: str, chunk_size: int = 128) -> List[Dict[str, Any]]: 
    chunks = get_chunks(text)
    print (text_list)
    chunk_embeddings = np.array(get_openai_embeddings(chunks)) 
    resp = add_vectors(faiss_client, chunk_embeddings)
    print(resp)  
    return "Uploaded Successfully"


def split_full_text(text, chunk_size=250):
    try:
        chunks = get_chunks(text, chunk_size=chunk_size)
        response = [{"text": chunk} for chunk in chunks]
        return response
    except Exception as e:
        print("*****************ALARM at split_full_text*****************")
        return str(e)

def get_openai_embeddings(texts, model="text-embedding-3-small"):
    try:
        texts = [text.replace("\n", " ") for text in texts]
        response = client.embeddings.create(input=texts, model=model)
        embeddings = [data.embedding for data in response.data]
        return embeddings
    except Exception as e:
        print ("*****************ALARM at open ai embeddings*****************")
        return str(e)

def get_openai_response(question, text):
    try:
        user_query = f"QUESTION: {question}\nTEXT: {text}"
        response_text = (client.chat.completions.create(
        model = 'gpt-4o-mini',
        messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_query}
        ]
        ).choices[0].message.content)
        if '\n' in response_text:
            print (response_text)
        if response_text[response_text.find('***')+3:response_text.find('***', response_text.find('***')+3)].lower() in text.lower():
           return response_text
        else:
            # text = text.replace('***', '')
            # highlighted_text = longest_common_substring(text, user_query)
            return "None"
        
    except Exception as e:
        print ("*****************ALARM at open ai repsonse*****************")
        return str(e)

